#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Скрипт для создания презентации в PowerPoint для Модуля 2:
Архитектура эффективных промптов
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Создаем презентацию
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Слайд 1: Титульный
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой слайд
    title1 = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title1.text_frame
    title_frame.text = "АРХИТЕКТУРА ЭФФЕКТИВНЫХ ПРОМПТОВ"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(41, 51, 92)
    title_para.alignment = PP_ALIGN.CENTER

    subtitle1 = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle1.text_frame
    subtitle_frame.text = "Модуль 2 | От хаоса к системе"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Слайд 2: 4 компонента
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Заголовок и контент
    title2 = slide2.shapes.title
    title2.text = "4 КОМПОНЕНТА ЭФФЕКТИВНОГО ПРОМПТА"

    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.text = "1. РОЛЬ (Who)"
    p = tf.paragraphs[0]
    p.level = 0
    p.font.size = Pt(24)
    
    tf.add_paragraph()
    tf.paragraphs[1].text = "   └ Кто ты? Эксперт по..."
    tf.paragraphs[1].level = 1
    tf.paragraphs[1].font.size = Pt(20)

    tf.add_paragraph()
    tf.paragraphs[2].text = ""
    tf.paragraphs[2].level = 0

    tf.add_paragraph()
    tf.paragraphs[3].text = "2. КОНТЕКСТ (What)"
    tf.paragraphs[3].level = 0
    tf.paragraphs[3].font.size = Pt(24)

    tf.add_paragraph()
    tf.paragraphs[4].text = "   └ Аудитория, цель, тон"
    tf.paragraphs[4].level = 1
    tf.paragraphs[4].font.size = Pt(20)

    tf.add_paragraph()
    tf.paragraphs[5].text = ""

    tf.add_paragraph()
    tf.paragraphs[6].text = "3. ИНСТРУКЦИЯ (How)"
    tf.paragraphs[6].level = 0
    tf.paragraphs[6].font.size = Pt(24)

    tf.add_paragraph()
    tf.paragraphs[7].text = "   └ Что конкретно сделать?"
    tf.paragraphs[7].level = 1
    tf.paragraphs[7].font.size = Pt(20)

    tf.add_paragraph()
    tf.paragraphs[8].text = ""

    tf.add_paragraph()
    tf.paragraphs[9].text = "4. ФОРМАТ (Output)"
    tf.paragraphs[9].level = 0
    tf.paragraphs[9].font.size = Pt(24)

    tf.add_paragraph()
    tf.paragraphs[10].text = "   └ JSON, Markdown, список..."
    tf.paragraphs[10].level = 1
    tf.paragraphs[10].font.size = Pt(20)

    # Слайд 3: Плохой vs Хороший пример
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "ПЛОХОЙ vs ХОРОШИЙ ПРОМПТ"

    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "❌ ПЛОХОЙ:"
    tf3.paragraphs[0].font.size = Pt(28)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)

    tf3.add_paragraph()
    tf3.paragraphs[1].text = '   "Напиши статью про маркетинг"'
    tf3.paragraphs[1].font.size = Pt(20)
    tf3.paragraphs[1].level = 1

    tf3.add_paragraph()
    tf3.paragraphs[2].text = ""

    tf3.add_paragraph()
    tf3.paragraphs[3].text = "✅ ХОРОШИЙ:"
    tf3.paragraphs[3].font.size = Pt(28)
    tf3.paragraphs[3].font.bold = True
    tf3.paragraphs[3].font.color.rgb = RGBColor(0, 150, 0)

    tf3.add_paragraph()
    tf3.paragraphs[4].text = "   Ты – эксперт по контент-маркетингу"
    tf3.paragraphs[4].font.size = Pt(18)
    tf3.paragraphs[4].level = 1

    tf3.add_paragraph()
    tf3.paragraphs[5].text = "   Контекст: B2B, маркетологи, вебинар"
    tf3.paragraphs[5].font.size = Pt(18)
    tf3.paragraphs[5].level = 1

    tf3.add_paragraph()
    tf3.paragraphs[6].text = "   Задача: статья на 1000 слов"
    tf3.paragraphs[6].font.size = Pt(18)
    tf3.paragraphs[6].level = 1

    tf3.add_paragraph()
    tf3.paragraphs[7].text = "   Формат: введение + 5 пунктов + CTA"
    tf3.paragraphs[7].font.size = Pt(18)
    tf3.paragraphs[7].level = 1

    # Слайд 4: 4 фреймворка (обзор)
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "4 ФРЕЙМВОРКА ПРОМПТИНГА"

    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "RISEN  → Бизнес-анализ, стратегия"
    tf4.paragraphs[0].font.size = Pt(24)
    
    tf4.add_paragraph()
    tf4.paragraphs[1].text = ""
    
    tf4.add_paragraph()
    tf4.paragraphs[2].text = "CRISPE → Креатив, идеи, слоганы"
    tf4.paragraphs[2].font.size = Pt(24)
    
    tf4.add_paragraph()
    tf4.paragraphs[3].text = ""
    
    tf4.add_paragraph()
    tf4.paragraphs[4].text = "CREATE → Контент, статьи, гайды"
    tf4.paragraphs[4].font.size = Pt(24)
    
    tf4.add_paragraph()
    tf4.paragraphs[5].text = ""
    
    tf4.add_paragraph()
    tf4.paragraphs[6].text = "RTF    → Быстрые задачи"
    tf4.paragraphs[6].font.size = Pt(24)

    # Слайд 5: RISEN
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "RISEN - для бизнес-задач"

    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.text = "R – Role (Роль)"
    tf5.paragraphs[0].font.size = Pt(22)
    
    tf5.add_paragraph()
    tf5.paragraphs[1].text = "I – Input (Данные)"
    tf5.paragraphs[1].font.size = Pt(22)
    
    tf5.add_paragraph()
    tf5.paragraphs[2].text = "S – Steps (Шаги)"
    tf5.paragraphs[2].font.size = Pt(22)
    
    tf5.add_paragraph()
    tf5.paragraphs[3].text = "E – Expectation (Результат)"
    tf5.paragraphs[3].font.size = Pt(22)
    
    tf5.add_paragraph()
    tf5.paragraphs[4].text = "N – Narrowing (Ограничения)"
    tf5.paragraphs[4].font.size = Pt(22)

    # Слайд 6: CRISPE
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "CRISPE - для креатива"

    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.text = "C – Capacity/Role (Роль)"
    tf6.paragraphs[0].font.size = Pt(20)
    
    tf6.add_paragraph()
    tf6.paragraphs[1].text = "R – Insight (Инсайт)"
    tf6.paragraphs[1].font.size = Pt(20)
    
    tf6.add_paragraph()
    tf6.paragraphs[2].text = "I – Statement (Утверждение)"
    tf6.paragraphs[2].font.size = Pt(20)
    
    tf6.add_paragraph()
    tf6.paragraphs[3].text = "S – Personality (Стиль)"
    tf6.paragraphs[3].font.size = Pt(20)
    
    tf6.add_paragraph()
    tf6.paragraphs[4].text = "P – Experiment (Варианты)"
    tf6.paragraphs[4].font.size = Pt(20)
    
    tf6.add_paragraph()
    tf6.paragraphs[5].text = "E – Expectation (Результат)"
    tf6.paragraphs[5].font.size = Pt(20)

    # Слайд 7: CREATE
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    title7.text = "CREATE - для контента"

    content7 = slide7.placeholders[1]
    tf7 = content7.text_frame
    tf7.text = "C – Character (Персонаж)"
    tf7.paragraphs[0].font.size = Pt(20)
    
    tf7.add_paragraph()
    tf7.paragraphs[1].text = "R – Request (Запрос)"
    tf7.paragraphs[1].font.size = Pt(20)
    
    tf7.add_paragraph()
    tf7.paragraphs[2].text = "E – Examples (Примеры)"
    tf7.paragraphs[2].font.size = Pt(20)
    
    tf7.add_paragraph()
    tf7.paragraphs[3].text = "A – Adjustments (Стиль)"
    tf7.paragraphs[3].font.size = Pt(20)
    
    tf7.add_paragraph()
    tf7.paragraphs[4].text = "T – Type (Формат)"
    tf7.paragraphs[4].font.size = Pt(20)
    
    tf7.add_paragraph()
    tf7.paragraphs[5].text = "E – Extras (Дополнительно)"
    tf7.paragraphs[5].font.size = Pt(20)

    # Слайд 8: RTF
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    title8.text = "RTF - для быстрых задач"

    content8 = slide8.placeholders[1]
    tf8 = content8.text_frame
    tf8.text = "R – Role (Роль)"
    tf8.paragraphs[0].font.size = Pt(28)
    
    tf8.add_paragraph()
    tf8.paragraphs[1].text = ""
    
    tf8.add_paragraph()
    tf8.paragraphs[2].text = "T – Task (Задача)"
    tf8.paragraphs[2].font.size = Pt(28)
    
    tf8.add_paragraph()
    tf8.paragraphs[3].text = ""
    
    tf8.add_paragraph()
    tf8.paragraphs[4].text = "F – Format (Формат)"
    tf8.paragraphs[4].font.size = Pt(28)

    # Слайд 9: Когда использовать
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    title9.text = "КАКОЙ ФРЕЙМВОРК ВЫБРАТЬ?"

    content9 = slide9.placeholders[1]
    tf9 = content9.text_frame
    tf9.text = "RISEN   → Анализ, стратегия, данные"
    tf9.paragraphs[0].font.size = Pt(22)
    
    tf9.add_paragraph()
    tf9.paragraphs[1].text = ""
    
    tf9.add_paragraph()
    tf9.paragraphs[2].text = "CRISPE  → Креатив, идеи, слоганы"
    tf9.paragraphs[2].font.size = Pt(22)
    
    tf9.add_paragraph()
    tf9.paragraphs[3].text = ""
    
    tf9.add_paragraph()
    tf9.paragraphs[4].text = "CREATE  → Статьи, гайды, контент"
    tf9.paragraphs[4].font.size = Pt(22)
    
    tf9.add_paragraph()
    tf9.paragraphs[5].text = ""
    
    tf9.add_paragraph()
    tf9.paragraphs[6].text = "RTF     → Быстрые запросы"
    tf9.paragraphs[6].font.size = Pt(22)

    # Слайд 10: Токены и лимиты
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    title10.text = "ТОКЕНЫ И ЛИМИТЫ"

    content10 = slide10.placeholders[1]
    tf10 = content10.text_frame
    tf10.text = "1 токен ≈ 0.75 слова (англ)"
    tf10.paragraphs[0].font.size = Pt(24)
    
    tf10.add_paragraph()
    tf10.paragraphs[1].text = "1 токен ≈ 0.5 слова (рус)"
    tf10.paragraphs[1].font.size = Pt(24)
    
    tf10.add_paragraph()
    tf10.paragraphs[2].text = ""
    tf10.paragraphs[2].font.bold = True
    tf10.paragraphs[2].font.size = Pt(22)
    tf10.paragraphs[2].text = "ChatGPT: 40 сообщений / 3 часа"
    
    tf10.add_paragraph()
    tf10.paragraphs[3].text = "Claude: ~50 сообщений / 4 часа"
    tf10.paragraphs[3].font.size = Pt(20)
    
    tf10.add_paragraph()
    tf10.paragraphs[4].text = ""
    
    tf10.add_paragraph()
    tf10.paragraphs[5].text = "💡 Совет: Пишите короче и точнее"
    tf10.paragraphs[5].font.size = Pt(24)
    tf10.paragraphs[5].font.italic = True

    # Слайд 11: Структурированный вывод
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    title11.text = "СТРУКТУРИРОВАННЫЙ ВЫВОД"

    content11 = slide11.placeholders[1]
    tf11 = content11.text_frame
    tf11.text = "JSON  → Данные, API"
    tf11.paragraphs[0].font.size = Pt(24)
    
    tf11.add_paragraph()
    tf11.paragraphs[1].text = "Markdown → Документы"
    tf11.paragraphs[1].font.size = Pt(24)
    
    tf11.add_paragraph()
    tf11.paragraphs[2].text = "CSV  → Таблицы"
    tf11.paragraphs[2].font.size = Pt(24)
    
    tf11.add_paragraph()
    tf11.paragraphs[3].text = "XML  → Конфигурации"
    tf11.paragraphs[3].font.size = Pt(24)
    
    tf11.add_paragraph()
    tf11.paragraphs[4].text = ""
    
    tf11.add_paragraph()
    tf11.paragraphs[5].text = "💡 Всегда показывайте пример структуры!"
    tf11.paragraphs[5].font.size = Pt(20)
    tf11.paragraphs[5].font.italic = True

    # Слайд 12: 7 правил
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    title12 = slide12.shapes.title
    title12.text = "7 ПРАВИЛ ЭФФЕКТИВНЫХ ПРОМПТОВ"

    content12 = slide12.placeholders[1]
    tf12 = content12.text_frame
    tf12.text = "1. Будьте специфичны"
    tf12.paragraphs[0].font.size = Pt(18)
    
    tf12.add_paragraph()
    tf12.paragraphs[1].text = "2. Используйте примеры"
    tf12.paragraphs[1].font.size = Pt(18)
    
    tf12.add_paragraph()
    tf12.paragraphs[2].text = "3. Разбивайте сложное"
    tf12.paragraphs[2].font.size = Pt(18)
    
    tf12.add_paragraph()
    tf12.paragraphs[3].text = "4. Указывайте формат"
    tf12.paragraphs[3].font.size = Pt(18)
    
    tf12.add_paragraph()
    tf12.paragraphs[4].text = "5. Итерируйте (улучшайте)"
    tf12.paragraphs[4].font.size = Pt(18)
    
    tf12.add_paragraph()
    tf12.paragraphs[5].text = "6. Тестируйте на разных моделях"
    tf12.paragraphs[5].font.size = Pt(18)
    
    tf12.add_paragraph()
    tf12.paragraphs[6].text = "7. Создавайте библиотеку шаблонов"
    tf12.paragraphs[6].font.size = Pt(18)

    # Слайд 13: Переход к практике
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    title13 = slide13.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame13 = title13.text_frame
    title_frame13.text = "🎯 ПЕРЕХОДИМ К ПРАКТИКЕ"
    title_para13 = title_frame13.paragraphs[0]
    title_para13.font.size = Pt(44)
    title_para13.font.bold = True
    title_para13.font.color.rgb = RGBColor(41, 51, 92)
    title_para13.alignment = PP_ALIGN.CENTER

    content13 = slide13.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(3))
    content_frame13 = content13.text_frame
    content_frame13.text = "• Применять фреймворки (30 мин)"
    content_frame13.paragraphs[0].font.size = Pt(24)
    
    content_frame13.add_paragraph()
    content_frame13.paragraphs[1].text = "• Структурированный вывод (25 мин)"
    content_frame13.paragraphs[1].font.size = Pt(24)
    
    content_frame13.add_paragraph()
    content_frame13.paragraphs[2].text = ""
    
    content_frame13.add_paragraph()
    content_frame13.paragraphs[3].text = "Минимум 7 промптов на занятии!"
    content_frame13.paragraphs[3].font.size = Pt(28)
    content_frame13.paragraphs[3].font.bold = True

    # Сохраняем презентацию
    output_file = "Модуль-2-Архитектура-промптов.pptx"
    prs.save(output_file)
    print(f"✅ Презентация создана: {output_file}")
    print(f"📊 Слайдов: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
