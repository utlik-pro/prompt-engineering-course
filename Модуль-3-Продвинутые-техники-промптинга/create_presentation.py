#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Скрипт для создания презентации в PowerPoint для Модуля 3:
Продвинутые техники промптинга
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Создаем презентацию
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Слайд 1: Титульный
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой слайд
    title1 = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title1.text_frame
    title_frame.text = "ПРОДВИНУТЫЕ ТЕХНИКИ\nПРОМПТИНГА"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(41, 51, 92)
    title_para.alignment = PP_ALIGN.CENTER

    subtitle1 = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    subtitle_frame = subtitle1.text_frame
    subtitle_frame.text = "Модуль 3 | Заставляем модель думать вслух"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Слайд 2: Проблема обычного промптинга
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "ПРОБЛЕМА ОБЫЧНОГО ПРОМПТИНГА"

    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.text = "❌ ОБЫЧНЫЙ ПРОМПТ:"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)
    
    tf.add_paragraph()
    tf.paragraphs[1].text = '   "Реши: У Маши было 5 яблок.'
    tf.paragraphs[1].font.size = Pt(18)
    tf.paragraphs[1].level = 1
    
    tf.add_paragraph()
    tf.paragraphs[2].text = "   Она съела 2. Сколько осталось?" + '"'
    tf.paragraphs[2].font.size = Pt(18)
    tf.paragraphs[2].level = 1

    tf.add_paragraph()
    tf.paragraphs[3].text = ""
    
    tf.add_paragraph()
    tf.paragraphs[4].text = "Ответ: '3'"
    tf.paragraphs[4].font.size = Pt(20)
    tf.paragraphs[4].level = 1

    tf.add_paragraph()
    tf.paragraphs[5].text = ""
    
    tf.add_paragraph()
    tf.paragraphs[6].text = "Проблема: Нет рассуждения, не видно логики"
    tf.paragraphs[6].font.size = Pt(20)
    tf.paragraphs[6].font.italic = True
    tf.paragraphs[6].font.color.rgb = RGBColor(150, 0, 0)

    # Слайд 3: Chain-of-Thought (CoT) - Основы
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "CHAIN-OF-THOUGHT (COT)"

    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Идея: Заставить модель показывать процесс рассуждения пошагово"
    tf3.paragraphs[0].font.size = Pt(22)
    
    tf3.add_paragraph()
    tf3.paragraphs[1].text = ""
    
    tf3.add_paragraph()
    tf3.paragraphs[2].text = "Формула:"
    tf3.paragraphs[2].font.size = Pt(24)
    tf3.paragraphs[2].font.bold = True
    
    tf3.add_paragraph()
    tf3.paragraphs[3].text = "   Задача + 'Давайте рассуждать пошагово'"
    tf3.paragraphs[3].font.size = Pt(20)
    tf3.paragraphs[3].level = 1
    
    tf3.add_paragraph()
    tf3.paragraphs[4].text = "   = Пошаговое решение"
    tf3.paragraphs[4].font.size = Pt(20)
    tf3.paragraphs[4].level = 1

    tf3.add_paragraph()
    tf3.paragraphs[5].text = ""
    
    tf3.add_paragraph()
    tf3.paragraphs[6].text = "Результаты точнее на 20-30%"
    tf3.paragraphs[6].font.size = Pt(20)
    tf3.paragraphs[6].font.color.rgb = RGBColor(0, 150, 0)
    tf3.paragraphs[6].font.bold = True

    # Слайд 4: Пример CoT
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "ПРИМЕР COT"

    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "✅ ПРОМПТ С COT:"
    tf4.paragraphs[0].font.size = Pt(22)
    tf4.paragraphs[0].font.bold = True
    tf4.paragraphs[0].font.color.rgb = RGBColor(0, 150, 0)
    
    tf4.add_paragraph()
    tf4.paragraphs[1].text = '   "Реши, рассуждая пошагово:'
    tf4.paragraphs[1].font.size = Pt(17)
    tf4.paragraphs[1].level = 1
    
    tf4.add_paragraph()
    tf4.paragraphs[2].text = "   У Маши было 5 яблок..."
    tf4.paragraphs[2].font.size = Pt(17)
    tf4.paragraphs[2].level = 1

    tf4.add_paragraph()
    tf4.paragraphs[3].text = ""
    
    tf4.add_paragraph()
    tf4.paragraphs[4].text = "✅ ОТВЕТ:"
    tf4.paragraphs[4].font.size = Pt(22)
    tf4.paragraphs[4].font.bold = True
    tf4.paragraphs[4].font.color.rgb = RGBColor(0, 150, 0)
    
    tf4.add_paragraph()
    tf4.paragraphs[5].text = "   Шаг 1: 5 яблок"
    tf4.paragraphs[5].font.size = Pt(18)
    tf4.paragraphs[5].level = 1
    
    tf4.add_paragraph()
    tf4.paragraphs[6].text = "   Шаг 2: Съела 2"
    tf4.paragraphs[6].font.size = Pt(18)
    tf4.paragraphs[6].level = 1
    
    tf4.add_paragraph()
    tf4.paragraphs[7].text = "   Шаг 3: 5 - 2 = 3"
    tf4.paragraphs[7].font.size = Pt(18)
    tf4.paragraphs[7].level = 1

    # Слайд 5: CoT для сложных задач
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "COT ДЛЯ СЛОЖНЫХ ЗАДАЧ"

    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.text = "Для сложных задач CoT необходим!"
    tf5.paragraphs[0].font.size = Pt(24)
    tf5.paragraphs[0].font.bold = True
    
    tf5.add_paragraph()
    tf5.paragraphs[1].text = ""
    
    tf5.add_paragraph()
    tf5.paragraphs[2].text = "Покажи:"
    tf5.paragraphs[2].font.size = Pt(22)
    
    tf5.add_paragraph()
    tf5.paragraphs[3].text = "1. Что известно"
    tf5.paragraphs[3].font.size = Pt(20)
    tf5.paragraphs[3].level = 1
    
    tf5.add_paragraph()
    tf5.paragraphs[4].text = "2. Что нужно найти"
    tf5.paragraphs[4].font.size = Pt(20)
    tf5.paragraphs[4].level = 1
    
    tf5.add_paragraph()
    tf5.paragraphs[5].text = "3. Каждый шаг решения"
    tf5.paragraphs[5].font.size = Pt(20)
    tf5.paragraphs[5].level = 1
    
    tf5.add_paragraph()
    tf5.paragraphs[6].text = "4. Проверку ответа"
    tf5.paragraphs[6].font.size = Pt(20)
    tf5.paragraphs[6].level = 1

    # Слайд 6: Tree of Thoughts (ToT)
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "TREE OF THOUGHTS (TOT)"

    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.text = "Идея: Генерировать несколько вариантов решения, затем выбрать лучший"
    tf6.paragraphs[0].font.size = Pt(20)
    
    tf6.add_paragraph()
    tf6.paragraphs[1].text = ""
    
    tf6.add_paragraph()
    tf6.paragraphs[2].text = "Когда использовать:"
    tf6.paragraphs[2].font.size = Pt(22)
    tf6.paragraphs[2].font.bold = True
    
    tf6.add_paragraph()
    tf6.paragraphs[3].text = "• Творческие задачи"
    tf6.paragraphs[3].font.size = Pt(20)
    tf6.paragraphs[3].level = 1
    
    tf6.add_paragraph()
    tf6.paragraphs[4].text = "• Задачи с множественными решениями"
    tf6.paragraphs[4].font.size = Pt(20)
    tf6.paragraphs[4].level = 1
    
    tf6.add_paragraph()
    tf6.paragraphs[5].text = "• Когда нет одного 'правильного' ответа"
    tf6.paragraphs[5].font.size = Pt(20)
    tf6.paragraphs[5].level = 1

    # Слайд 7: Self-Consistency
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    title7.text = "SELF-CONSISTENCY"

    content7 = slide7.placeholders[1]
    tf7 = content7.text_frame
    tf7.text = "Идея: Запустить промпт несколько раз и выбрать наиболее частый ответ"
    tf7.paragraphs[0].font.size = Pt(20)
    
    tf7.add_paragraph()
    tf7.paragraphs[1].text = ""
    
    tf7.add_paragraph()
    tf7.paragraphs[2].text = "Метод:"
    tf7.paragraphs[2].font.size = Pt(22)
    tf7.paragraphs[2].font.bold = True
    
    tf7.add_paragraph()
    tf7.paragraphs[3].text = "1. Запустить промпт 3-5 раз"
    tf7.paragraphs[3].font.size = Pt(20)
    tf7.paragraphs[3].level = 1
    
    tf7.add_paragraph()
    tf7.paragraphs[4].text = "2. Собрать ответы"
    tf7.paragraphs[4].font.size = Pt(20)
    tf7.paragraphs[4].level = 1
    
    tf7.add_paragraph()
    tf7.paragraphs[5].text = "3. Выбрать наиболее частый ответ (консенсус)"
    tf7.paragraphs[5].font.size = Pt(20)
    tf7.paragraphs[5].level = 1

    tf7.add_paragraph()
    tf7.paragraphs[6].text = ""
    
    tf7.add_paragraph()
    tf7.paragraphs[7].text = "Когда: Критически важные задачи, нужна точность"
    tf7.paragraphs[7].font.size = Pt(18)
    tf7.paragraphs[7].font.italic = True

    # Слайд 8: ReAct (Reasoning + Acting)
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    title8.text = "REACT (REASONING + ACTING)"

    content8 = slide8.placeholders[1]
    tf8 = content8.text_frame
    tf8.text = "Идея: Чередовать рассуждение и действие в цикле"
    tf8.paragraphs[0].font.size = Pt(22)
    
    tf8.add_paragraph()
    tf8.paragraphs[1].text = ""
    
    tf8.add_paragraph()
    tf8.paragraphs[2].text = "Структура:"
    tf8.paragraphs[2].font.size = Pt(22)
    tf8.paragraphs[2].font.bold = True
    
    tf8.add_paragraph()
    tf8.paragraphs[3].text = "1. Thought (Думаю) - анализ ситуации"
    tf8.paragraphs[3].font.size = Pt(20)
    tf8.paragraphs[3].level = 1
    
    tf8.add_paragraph()
    tf8.paragraphs[4].text = "2. Action (Действую) - конкретное действие"
    tf8.paragraphs[4].font.size = Pt(20)
    tf8.paragraphs[4].level = 1
    
    tf8.add_paragraph()
    tf8.paragraphs[5].text = "3. Observation (Наблюдаю) - результат"
    tf8.paragraphs[5].font.size = Pt(20)
    tf8.paragraphs[5].level = 1
    
    tf8.add_paragraph()
    tf8.paragraphs[6].text = "4. Повтор до решения"
    tf8.paragraphs[6].font.size = Pt(20)
    tf8.paragraphs[6].level = 1

    # Слайд 9: Meta-prompting и рефлексия
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    title9.text = "META-PROMPTING"

    content9 = slide9.placeholders[1]
    tf9 = content9.text_frame
    tf9.text = "Идея: Промпт, который создает промпты"
    tf9.paragraphs[0].font.size = Pt(22)
    
    tf9.add_paragraph()
    tf9.paragraphs[1].text = ""
    
    tf9.add_paragraph()
    tf9.paragraphs[2].text = "Пример:"
    tf9.paragraphs[2].font.size = Pt(22)
    tf9.paragraphs[2].font.bold = True
    
    tf9.add_paragraph()
    tf9.paragraphs[3].text = '   "Ты – эксперт по промптам.'
    tf9.paragraphs[3].font.size = Pt(18)
    tf9.paragraphs[3].level = 1
    
    tf9.add_paragraph()
    tf9.paragraphs[4].text = "   Создай промпт для задачи X,"
    tf9.paragraphs[4].font.size = Pt(18)
    tf9.paragraphs[4].level = 1
    
    tf9.add_paragraph()
    tf9.paragraphs[5].text = "   используя CoT и JSON-вывод" + '"'
    tf9.paragraphs[5].font.size = Pt(18)
    tf9.paragraphs[5].level = 1

    tf9.add_paragraph()
    tf9.paragraphs[6].text = ""
    
    tf9.add_paragraph()
    tf9.paragraphs[7].text = "Рефлексия: Модель анализирует свой ответ и улучшает его"
    tf9.paragraphs[7].font.size = Pt(18)
    tf9.paragraphs[7].font.italic = True

    # Слайд 10: Сравнение техник
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    title10.text = "КОГДА ИСПОЛЬЗОВАТЬ КАКУЮ ТЕХНИКУ?"

    content10 = slide10.placeholders[1]
    tf10 = content10.text_frame
    tf10.text = "CoT → Сложные задачи с вычислениями"
    tf10.paragraphs[0].font.size = Pt(20)
    
    tf10.add_paragraph()
    tf10.paragraphs[1].text = ""
    
    tf10.add_paragraph()
    tf10.paragraphs[2].text = "ToT → Творческие задачи, множественные варианты"
    tf10.paragraphs[2].font.size = Pt(20)
    
    tf10.add_paragraph()
    tf10.paragraphs[3].text = ""
    
    tf10.add_paragraph()
    tf10.paragraphs[4].text = "Self-Consistency → Критически важные задачи"
    tf10.paragraphs[4].font.size = Pt(20)
    
    tf10.add_paragraph()
    tf10.paragraphs[5].text = ""
    
    tf10.add_paragraph()
    tf10.paragraphs[6].text = "ReAct → Динамические задачи, итеративные действия"
    tf10.paragraphs[6].font.size = Pt(20)
    
    tf10.add_paragraph()
    tf10.paragraphs[7].text = ""
    
    tf10.add_paragraph()
    tf10.paragraphs[8].text = "Meta-prompting → Создание промптов автоматически"
    tf10.paragraphs[8].font.size = Pt(20)

    # Слайд 11: Комбинирование техник
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    title11.text = "КОМБИНИРОВАНИЕ ТЕХНИК"

    content11 = slide11.placeholders[1]
    tf11 = content11.text_frame
    tf11.text = "Пример 1: CoT + Self-Consistency"
    tf11.paragraphs[0].font.size = Pt(22)
    tf11.paragraphs[0].font.bold = True
    
    tf11.add_paragraph()
    tf11.paragraphs[1].text = "1. Используйте CoT для пошагового решения"
    tf11.paragraphs[1].font.size = Pt(18)
    tf11.paragraphs[1].level = 1
    
    tf11.add_paragraph()
    tf11.paragraphs[2].text = "2. Запустите несколько раз"
    tf11.paragraphs[2].font.size = Pt(18)
    tf11.paragraphs[2].level = 1
    
    tf11.add_paragraph()
    tf11.paragraphs[3].text = "3. Выберите наиболее логичное решение"
    tf11.paragraphs[3].font.size = Pt(18)
    tf11.paragraphs[3].level = 1

    tf11.add_paragraph()
    tf11.paragraphs[4].text = ""
    
    tf11.add_paragraph()
    tf11.paragraphs[5].text = "Пример 2: ToT + CoT"
    tf11.paragraphs[5].font.size = Pt(22)
    tf11.paragraphs[5].font.bold = True
    
    tf11.add_paragraph()
    tf11.paragraphs[6].text = "1. ToT генерирует варианты"
    tf11.paragraphs[6].font.size = Pt(18)
    tf11.paragraphs[6].level = 1
    
    tf11.add_paragraph()
    tf11.paragraphs[7].text = "2. CoT анализирует каждый вариант"
    tf11.paragraphs[7].font.size = Pt(18)
    tf11.paragraphs[7].level = 1

    # Слайд 12: Практические советы
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    title12 = slide12.shapes.title
    title12.text = "7 ПРАВИЛ ПРОДВИНУТОГО ПРОМПТИНГА"

    content12 = slide12.placeholders[1]
    tf12 = content12.text_frame
    tf12.text = "1. Всегда используйте CoT для сложных задач"
    tf12.paragraphs[0].font.size = Pt(18)
    
    tf12.add_paragraph()
    tf12.paragraphs[1].text = "2. Проверяйте рассуждение модели"
    tf12.paragraphs[1].font.size = Pt(18)
    
    tf12.add_paragraph()
    tf12.paragraphs[2].text = "3. Для критичных задач - Self-Consistency (3-5 запусков)"
    tf12.paragraphs[2].font.size = Pt(18)
    
    tf12.add_paragraph()
    tf12.paragraphs[3].text = "4. ToT - для задач без единственного правильного ответа"
    tf12.paragraphs[3].font.size = Pt(18)
    
    tf12.add_paragraph()
    tf12.paragraphs[4].text = "5. ReAct - для динамических задач"
    tf12.paragraphs[4].font.size = Pt(18)
    
    tf12.add_paragraph()
    tf12.paragraphs[5].text = "6. Комбинируйте техники для максимального эффекта"
    tf12.paragraphs[5].font.size = Pt(18)
    
    tf12.add_paragraph()
    tf12.paragraphs[6].text = "7. Сохраняйте успешные промпты с продвинутыми техниками"
    tf12.paragraphs[6].font.size = Pt(18)

    # Слайд 13: Переход к практике
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    title13 = slide13.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame13 = title13.text_frame
    title_frame13.text = "🎯 ПЕРЕХОДИМ К ПРАКТИКЕ"
    title_para13 = title_frame13.paragraphs[0]
    title_para13.font.size = Pt(44)
    title_para13.font.bold = True
    title_para13.font.color.rgb = RGBColor(41, 51, 92)
    title_para13.alignment = PP_ALIGN.CENTER

    content13 = slide13.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(3))
    content_frame13 = content13.text_frame
    content_frame13.text = "1. CoT на математических задачах (30 мин)"
    content_frame13.paragraphs[0].font.size = Pt(22)
    
    content_frame13.add_paragraph()
    content_frame13.paragraphs[1].text = "2. ToT для творческих задач (30 мин)"
    content_frame13.paragraphs[1].font.size = Pt(22)
    
    content_frame13.add_paragraph()
    content_frame13.paragraphs[2].text = "3. Self-Consistency (15 мин)"
    content_frame13.paragraphs[2].font.size = Pt(22)
    
    content_frame13.add_paragraph()
    content_frame13.paragraphs[3].text = "4. Создание 'мыслящих' промптов (15 мин)"
    content_frame13.paragraphs[3].font.size = Pt(22)

    # Слайд 14: Итоги теории
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    title14 = slide14.shapes.title
    title14.text = "ЧТО МЫ УЗНАЛИ"

    content14 = slide14.placeholders[1]
    tf14 = content14.text_frame
    tf14.text = "✓ Chain-of-Thought (CoT) → Пошаговое рассуждение"
    tf14.paragraphs[0].font.size = Pt(20)
    
    tf14.add_paragraph()
    tf14.paragraphs[1].text = ""
    
    tf14.add_paragraph()
    tf14.paragraphs[2].text = "✓ Tree of Thoughts (ToT) → Множественные варианты"
    tf14.paragraphs[2].font.size = Pt(20)
    
    tf14.add_paragraph()
    tf14.paragraphs[3].text = ""
    
    tf14.add_paragraph()
    tf14.paragraphs[4].text = "✓ Self-Consistency → Повышение точности"
    tf14.paragraphs[4].font.size = Pt(20)
    
    tf14.add_paragraph()
    tf14.paragraphs[5].text = ""
    
    tf14.add_paragraph()
    tf14.paragraphs[6].text = "✓ ReAct → Рассуждение + действие"
    tf14.paragraphs[6].font.size = Pt(20)
    
    tf14.add_paragraph()
    tf14.paragraphs[7].text = ""
    
    tf14.add_paragraph()
    tf14.paragraphs[8].text = "✓ Meta-prompting → Создание промптов промптами"
    tf14.paragraphs[8].font.size = Pt(20)

    # Сохраняем презентацию
    output_file = "Модуль-3-Продвинутые-техники-промптинга.pptx"
    prs.save(output_file)
    print(f"✅ Презентация создана: {output_file}")
    print(f"📊 Слайдов: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()

